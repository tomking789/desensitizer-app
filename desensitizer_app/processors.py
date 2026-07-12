from __future__ import annotations

import shutil
import re
from io import BytesIO
from collections import Counter
from pathlib import Path
from typing import Iterable

from .core import ZWSP_MARKER, _prefix_conflict_codes

import fitz
import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas

from .candidates import CandidateHit, ReplacementSpec
from .core import (
    DesensitizeError,
    SkippedFile,
    anonymize_text,
    apply_replacements_text,
    read_text_file,
    write_text_file,
)
from .mapping import MappingStore
from .rules import find_sensitive_spans


TEXT_SUFFIXES = {".txt", ".md", ".csv", ".log", ".json", ".xml", ".html", ".htm"}
WORD_SUFFIXES = {".docx"}
EXCEL_SUFFIXES = {".xlsx"}
POWERPOINT_SUFFIXES = {".pptx"}
PDF_SUFFIXES = {".pdf"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
LEGACY_OFFICE_SUFFIXES = {".doc", ".xls"}
PROCESSABLE_SUFFIXES = TEXT_SUFFIXES | WORD_SUFFIXES | EXCEL_SUFFIXES | POWERPOINT_SUFFIXES | PDF_SUFFIXES
KNOWN_SUFFIXES = PROCESSABLE_SUFFIXES | IMAGE_SUFFIXES | LEGACY_OFFICE_SUFFIXES
PDF_RESTORE_FONT = "STSong-Light"
PDF_HEADER_FOOTER_REGION_RATIO = 0.08
PDF_HEADER_FOOTER_MIN_PAGES = 2
EXCEL_HEADER_SCAN_ROWS = 15
PERSON_NAME_HEADERS = {
    "姓名",
    "客户姓名",
    "员工姓名",
    "人员姓名",
    "成员姓名",
    "高管姓名",
    "联系人",
    "负责人",
    "法定代表人",
    "法人",
    "经办人",
    "申请人",
    "收件人",
    "患者姓名",
    "当事人",
}
SUBJECT_ID_HEADERS = {
    "编号",
    "id",
    "ID",
    "客户编号",
    "客户号",
    "客户ID",
    "客户id",
    "员工编号",
    "员工号",
    "工号",
    "人员编号",
    "用户编号",
    "用户ID",
    "用户id",
    "成员编号",
    "学号",
    "证件号",
    "身份证号",
    "手机号",
    "手机",
    "电话",
    "邮箱",
    "邮件",
}


def anonymize_file(
    input_path: Path,
    output_dir: Path,
    mapping: MappingStore,
    custom_terms: Iterable[str] = (),
) -> tuple[Path, Counter[str], str]:
    suffix = input_path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return _anonymize_text_file(input_path, output_dir, mapping, custom_terms)
    if suffix in WORD_SUFFIXES:
        return _anonymize_docx(input_path, output_dir, mapping, custom_terms)
    if suffix in EXCEL_SUFFIXES:
        return _anonymize_xlsx(input_path, output_dir, mapping, custom_terms)
    if suffix in POWERPOINT_SUFFIXES:
        return _anonymize_pptx(input_path, output_dir, mapping, custom_terms)
    if suffix in PDF_SUFFIXES:
        return _anonymize_pdf_text(input_path, output_dir, mapping, custom_terms)
    if suffix in IMAGE_SUFFIXES:
        raise SkippedFile("图片文件需要 OCR。当前轻量版不做 OCR，请先转换为文字版 PDF、Word 或文本后再处理。")
    if suffix in LEGACY_OFFICE_SUFFIXES:
        raise SkippedFile("旧版 Office 格式请先用 Word/Excel/WPS 另存为 .docx 或 .xlsx 后再处理。")
    raise DesensitizeError(f"Unsupported file type: {suffix or '(no extension)'}")


def is_known_file(path: Path) -> bool:
    return path.suffix.lower() in KNOWN_SUFFIXES


def is_processable_file(path: Path) -> bool:
    return path.suffix.lower() in PROCESSABLE_SUFFIXES


def scan_file_candidates(
    input_path: Path,
    custom_terms: Iterable[str] = (),
) -> tuple[list[CandidateHit], str]:
    suffix = input_path.suffix.lower()
    if suffix in IMAGE_SUFFIXES:
        raise SkippedFile("图片文件需要 OCR。当前轻量版不做 OCR，请先转换为文字版 PDF、Word 或文本后再处理。")
    if suffix in LEGACY_OFFICE_SUFFIXES:
        raise SkippedFile("旧版 Office 格式请先用 Word/Excel/WPS 另存为 .docx 或 .xlsx 后再处理。")
    if suffix not in PROCESSABLE_SUFFIXES:
        raise DesensitizeError(f"Unsupported file type: {suffix or '(no extension)'}")
    if suffix in EXCEL_SUFFIXES:
        return _scan_xlsx_candidates(input_path, custom_terms)
    counts: Counter[tuple[str, str, str]] = Counter()
    for chunk in _extract_text_chunks(input_path):
        for finding in find_sensitive_spans(chunk, custom_terms):
            counts[(finding.entity, finding.prefix, finding.value)] += 1
    hits = [
        CandidateHit(
            entity=entity,
            prefix=prefix,
            value=value,
            count=count,
            file=input_path,
        )
        for (entity, prefix, value), count in counts.items()
    ]
    return hits, f"{len(hits)} unique candidate(s) found."


def _scan_xlsx_candidates(
    input_path: Path,
    custom_terms: Iterable[str] = (),
) -> tuple[list[CandidateHit], str]:
    general_counts: Counter[tuple[str, str, str]] = Counter()
    scoped_counts: Counter[tuple[str, str, str, str, str]] = Counter()
    workbook = load_workbook(input_path, data_only=False, read_only=True)
    try:
        for sheet in workbook.worksheets:
            header_row_index, headers = _detect_table_headers(sheet)
            name_columns = _columns_matching(headers, PERSON_NAME_HEADERS)
            id_columns = _columns_matching(headers, SUBJECT_ID_HEADERS)
            for row in sheet.iter_rows():
                row_values = {cell.column: _cell_text(cell.value) for cell in row}
                for value in row_values.values():
                    if not value:
                        continue
                    for finding in find_sensitive_spans(value, custom_terms):
                        general_counts[(finding.entity, finding.prefix, finding.value)] += 1

                if not header_row_index or row[0].row <= header_row_index or not name_columns or not id_columns:
                    continue
                identifier = _first_identifier(row_values, headers, id_columns)
                if not identifier:
                    continue
                id_header, id_value = identifier
                context_key = _subject_context_key("PERSON", id_header, id_value)
                context_label = f"{id_header}: {id_value}"
                for column_index in name_columns:
                    name = row_values.get(column_index, "")
                    if not _looks_like_person_name(name):
                        continue
                    scoped_counts[("PERSON", "PERSON", name, context_key, context_label)] += 1
    finally:
        workbook.close()

    hits = [
        CandidateHit(
            entity=entity,
            prefix=prefix,
            value=value,
            count=count,
            file=input_path,
        )
        for (entity, prefix, value), count in general_counts.items()
    ]
    hits.extend(
        CandidateHit(
            entity=entity,
            prefix=prefix,
            value=value,
            count=count,
            file=input_path,
            context_key=context_key,
            context_label=context_label,
            source="excel-entity",
        )
        for (entity, prefix, value, context_key, context_label), count in scoped_counts.items()
    )
    return hits, f"{len(hits)} unique candidate(s) found; {len(scoped_counts)} row-scoped entity candidate(s)."


def anonymize_file_with_replacements(
    input_path: Path,
    output_dir: Path,
    replacements: Iterable[ReplacementSpec],
    remove_headers_footers: bool = False,
) -> tuple[Path, Counter[str], str]:
    suffix = input_path.suffix.lower()
    replacement_list = list(replacements)
    if suffix in TEXT_SUFFIXES:
        return _anonymize_text_file_with_replacements(input_path, output_dir, replacement_list)
    if suffix in WORD_SUFFIXES:
        return _anonymize_docx_with_replacements(input_path, output_dir, replacement_list, remove_headers_footers)
    if suffix in EXCEL_SUFFIXES:
        return _anonymize_xlsx_with_replacements(input_path, output_dir, replacement_list)
    if suffix in POWERPOINT_SUFFIXES:
        return _anonymize_pptx_with_replacements(input_path, output_dir, replacement_list)
    if suffix in PDF_SUFFIXES:
        return _anonymize_pdf_text_with_replacements(input_path, output_dir, replacement_list, remove_headers_footers)
    if suffix in IMAGE_SUFFIXES:
        raise SkippedFile("图片文件需要 OCR。当前轻量版不做 OCR，请先转换为文字版 PDF、Word 或文本后再处理。")
    if suffix in LEGACY_OFFICE_SUFFIXES:
        raise SkippedFile("旧版 Office 格式请先用 Word/Excel/WPS 另存为 .docx 或 .xlsx 后再处理。")
    raise DesensitizeError(f"Unsupported file type: {suffix or '(no extension)'}")


def restore_file(input_path: Path, output_dir: Path, mapping: MappingStore) -> tuple[Path, str]:
    suffix = input_path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        text = read_text_file(input_path)
        output_path = _output_path(input_path, output_dir, "restored", suffix)
        write_text_file(output_path, mapping.restore_text(text))
        return output_path, "Text restored."
    if suffix in WORD_SUFFIXES:
        output_path = _output_path(input_path, output_dir, "restored", suffix)
        document = Document(str(input_path))
        _restore_docx_document(document, mapping)
        document.save(str(output_path))
        return output_path, "Word document restored."
    if suffix in EXCEL_SUFFIXES:
        output_path = _output_path(input_path, output_dir, "restored", suffix)
        workbook = load_workbook(input_path)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str):
                        cell.value = mapping.restore_text(cell.value)
        workbook.save(output_path)
        return output_path, "Excel workbook restored."
    if suffix in POWERPOINT_SUFFIXES:
        output_path = _output_path(input_path, output_dir, "restored", suffix)
        presentation = Presentation(str(input_path))
        _restore_pptx_deck(presentation, mapping)
        presentation.save(str(output_path))
        return output_path, "PowerPoint deck restored."
    if suffix in PDF_SUFFIXES:
        return _restore_pdf_text(input_path, output_dir, mapping)
    raise DesensitizeError("This file type cannot be restored by the current version.")


def remove_headers_footers(input_path: Path, output_dir: Path) -> tuple[Path, str]:
    suffix = input_path.suffix.lower()
    if suffix in WORD_SUFFIXES:
        return _remove_docx_headers_footers(input_path, output_dir)
    if suffix in PDF_SUFFIXES:
        return _remove_pdf_headers_footers(input_path, output_dir)
    raise SkippedFile("仅支持去除 .docx 和 .pdf 文件的页眉页脚；其他格式已跳过。")


def _remove_docx_headers_footers(input_path: Path, output_dir: Path) -> tuple[Path, str]:
    output_path = _output_path(input_path, output_dir, "no_header_footer", input_path.suffix)
    document = Document(str(input_path))
    cleared_parts = _clear_docx_headers_footers(document)
    document.save(str(output_path))
    return output_path, f"Removed header/footer content from {cleared_parts} part(s)."


def _remove_pdf_headers_footers(input_path: Path, output_dir: Path) -> tuple[Path, str]:
    output_path = _output_path(input_path, output_dir, "no_header_footer", input_path.suffix)
    document = fitz.open(str(input_path))
    try:
        removed_count = _apply_pdf_header_footer_redactions(document)
        if removed_count:
            document.save(str(output_path), garbage=4, deflate=True, clean=True)
        else:
            shutil.copy2(input_path, output_path)
    finally:
        document.close()
    if removed_count:
        return output_path, f"Removed {removed_count} repeated header/footer item(s)."
    return output_path, "No repeated PDF header/footer text was detected; file copied unchanged."


def _anonymize_text_file(
    input_path: Path,
    output_dir: Path,
    mapping: MappingStore,
    custom_terms: Iterable[str],
) -> tuple[Path, Counter[str], str]:
    text = read_text_file(input_path)
    new_text, counts = anonymize_text(text, mapping, custom_terms)
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    write_text_file(output_path, new_text)
    return output_path, counts, "Text file processed."


def _anonymize_docx(
    input_path: Path,
    output_dir: Path,
    mapping: MappingStore,
    custom_terms: Iterable[str],
) -> tuple[Path, Counter[str], str]:
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    document = Document(str(input_path))
    counts = Counter()
    counts.update(_replace_docx_document(document, mapping, custom_terms))
    document.save(str(output_path))
    return output_path, counts, "Word document processed. Text split across runs may need manual review."


def _anonymize_xlsx(
    input_path: Path,
    output_dir: Path,
    mapping: MappingStore,
    custom_terms: Iterable[str],
) -> tuple[Path, Counter[str], str]:
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    workbook = load_workbook(input_path)
    counts = Counter()
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if not isinstance(cell.value, str):
                    continue
                if cell.data_type == "f":
                    continue
                new_value, cell_counts = anonymize_text(cell.value, mapping, custom_terms)
                cell.value = new_value
                counts.update(cell_counts)
    workbook.save(output_path)
    return output_path, counts, "Excel workbook processed."


def _anonymize_pptx(
    input_path: Path,
    output_dir: Path,
    mapping: MappingStore,
    custom_terms: Iterable[str],
) -> tuple[Path, Counter[str], str]:
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    presentation = Presentation(str(input_path))
    counts = _replace_pptx_deck(presentation, mapping, custom_terms)
    presentation.save(str(output_path))
    return output_path, counts, "PowerPoint deck processed."


def _anonymize_pdf_text(
    input_path: Path,
    output_dir: Path,
    mapping: MappingStore,
    custom_terms: Iterable[str],
) -> tuple[Path, Counter[str], str]:
    replacements: list[ReplacementSpec] = []
    for chunk in _extract_pdf_text_chunks(input_path):
        for finding in find_sensitive_spans(chunk, custom_terms):
            placeholder = mapping.get_or_add(finding.value, finding.entity, finding.prefix)
            replacements.append(
                ReplacementSpec(
                    entity=finding.entity,
                    prefix=finding.prefix,
                    value=finding.value,
                    replacement=placeholder,
                )
            )
    if not replacements:
        output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
        shutil.copy2(input_path, output_path)
        return output_path, Counter(), "PDF contained extractable text but no approved sensitive value was found."
    return _redact_pdf_with_replacements(input_path, output_dir, replacements)


def _anonymize_text_file_with_replacements(
    input_path: Path,
    output_dir: Path,
    replacements: list[ReplacementSpec],
) -> tuple[Path, Counter[str], str]:
    text = read_text_file(input_path)
    new_text, counts = apply_replacements_text(text, replacements)
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    write_text_file(output_path, new_text)
    return output_path, counts, "Text file processed with approved replacements."


def _anonymize_docx_with_replacements(
    input_path: Path,
    output_dir: Path,
    replacements: list[ReplacementSpec],
    remove_headers_footers: bool = False,
) -> tuple[Path, Counter[str], str]:
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    document = Document(str(input_path))
    counts = Counter()
    cleared_parts = _clear_docx_headers_footers(document) if remove_headers_footers else 0
    counts.update(_apply_replacements_docx_document(document, replacements))
    document.save(str(output_path))
    message = "Word document processed with approved replacements."
    if remove_headers_footers:
        message += f" Removed header/footer content from {cleared_parts} part(s)."
    return output_path, counts, message


def _anonymize_xlsx_with_replacements(
    input_path: Path,
    output_dir: Path,
    replacements: list[ReplacementSpec],
) -> tuple[Path, Counter[str], str]:
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    workbook = load_workbook(input_path)
    counts = Counter()
    scoped_by_context = _replacements_by_context(replacements)
    global_replacements = [replacement for replacement in replacements if not replacement.context_key]
    for sheet in workbook.worksheets:
        header_row_index, headers = _detect_table_headers(sheet)
        id_columns = _columns_matching(headers, SUBJECT_ID_HEADERS)
        for row in sheet.iter_rows():
            row_values = {cell.column: _cell_text(cell.value) for cell in row}
            row_context_keys = _row_context_keys(row_values, headers, id_columns) if header_row_index and row[0].row > header_row_index else []
            scoped_replacements = _scoped_replacements_for_row(row_context_keys, scoped_by_context)
            for cell in row:
                if not isinstance(cell.value, str):
                    continue
                if cell.data_type == "f":
                    continue
                new_value, cell_counts = _apply_scoped_replacements_text(cell.value, scoped_replacements)
                if global_replacements:
                    new_value, global_counts = apply_replacements_text(new_value, global_replacements)
                    cell_counts.update(global_counts)
                cell.value = new_value
                counts.update(cell_counts)
    workbook.save(output_path)
    return output_path, counts, "Excel workbook processed with approved replacements."


def _anonymize_pptx_with_replacements(
    input_path: Path,
    output_dir: Path,
    replacements: list[ReplacementSpec],
) -> tuple[Path, Counter[str], str]:
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    presentation = Presentation(str(input_path))
    counts = _apply_replacements_pptx_deck(presentation, replacements)
    presentation.save(str(output_path))
    return output_path, counts, "PowerPoint deck processed with approved replacements."


def _anonymize_pdf_text_with_replacements(
    input_path: Path,
    output_dir: Path,
    replacements: list[ReplacementSpec],
    remove_headers_footers: bool = False,
) -> tuple[Path, Counter[str], str]:
    _ensure_pdf_has_extractable_text(input_path)
    return _redact_pdf_with_replacements(input_path, output_dir, replacements, remove_headers_footers)


def _replace_docx_document(
    document: Document,
    mapping: MappingStore,
    custom_terms: Iterable[str],
) -> Counter[str]:
    counts = Counter()
    for paragraph in document.paragraphs:
        counts.update(_replace_paragraph(paragraph, mapping, custom_terms))
    for table in document.tables:
        counts.update(_replace_docx_table(table, mapping, custom_terms))
    for section in document.sections:
        counts.update(_replace_header_footer(section.header, mapping, custom_terms))
        counts.update(_replace_header_footer(section.footer, mapping, custom_terms))
    return counts


def _replace_docx_table(table, mapping: MappingStore, custom_terms: Iterable[str]) -> Counter[str]:
    counts = Counter()
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                counts.update(_replace_paragraph(paragraph, mapping, custom_terms))
            for nested_table in cell.tables:
                counts.update(_replace_docx_table(nested_table, mapping, custom_terms))
    return counts


def _replace_header_footer(part, mapping: MappingStore, custom_terms: Iterable[str]) -> Counter[str]:
    counts = Counter()
    for paragraph in part.paragraphs:
        counts.update(_replace_paragraph(paragraph, mapping, custom_terms))
    for table in part.tables:
        counts.update(_replace_docx_table(table, mapping, custom_terms))
    return counts


def _apply_replacements_docx_document(
    document: Document,
    replacements: list[ReplacementSpec],
) -> Counter[str]:
    counts = Counter()
    for paragraph in document.paragraphs:
        counts.update(_apply_replacements_paragraph(paragraph, replacements))
    for table in document.tables:
        counts.update(_apply_replacements_docx_table(table, replacements))
    for section in document.sections:
        counts.update(_apply_replacements_header_footer(section.header, replacements))
        counts.update(_apply_replacements_header_footer(section.footer, replacements))
    return counts


def _apply_replacements_docx_table(table, replacements: list[ReplacementSpec]) -> Counter[str]:
    counts = Counter()
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                counts.update(_apply_replacements_paragraph(paragraph, replacements))
            for nested_table in cell.tables:
                counts.update(_apply_replacements_docx_table(nested_table, replacements))
    return counts


def _apply_replacements_header_footer(part, replacements: list[ReplacementSpec]) -> Counter[str]:
    counts = Counter()
    for paragraph in part.paragraphs:
        counts.update(_apply_replacements_paragraph(paragraph, replacements))
    for table in part.tables:
        counts.update(_apply_replacements_docx_table(table, replacements))
    return counts


def _replace_paragraph(paragraph, mapping: MappingStore, custom_terms: Iterable[str]) -> Counter[str]:
    if not paragraph.text:
        return Counter()
    new_text, counts = anonymize_text(paragraph.text, mapping, custom_terms)
    if counts:
        _set_paragraph_text(paragraph, new_text)
    return counts


def _apply_replacements_paragraph(paragraph, replacements: list[ReplacementSpec]) -> Counter[str]:
    if not paragraph.text:
        return Counter()
    new_text, counts = apply_replacements_text(paragraph.text, replacements)
    if counts:
        _set_paragraph_text(paragraph, new_text)
    return counts


def _set_paragraph_text(paragraph, text: str) -> None:
    # Replacing the whole paragraph catches sensitive text split across Word runs.
    # This may flatten mixed inline formatting in that paragraph, but prevents leaks.
    if not paragraph.runs:
        paragraph.add_run(text)
        return
    paragraph.runs[0].text = text
    for run in paragraph.runs[1:]:
        run.text = ""


def _apply_replacements_runs(runs, replacements: list[ReplacementSpec]) -> Counter[str]:
    counts = Counter()
    for run in runs:
        if not run.text:
            continue
        new_text, run_counts = apply_replacements_text(run.text, replacements)
        run.text = new_text
        counts.update(run_counts)
    return counts


def _replace_runs(runs, mapping: MappingStore, custom_terms: Iterable[str]) -> Counter[str]:
    counts = Counter()
    for run in runs:
        if not run.text:
            continue
        new_text, run_counts = anonymize_text(run.text, mapping, custom_terms)
        run.text = new_text
        counts.update(run_counts)
    return counts


def _restore_docx_document(document: Document, mapping: MappingStore) -> None:
    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            if run.text:
                run.text = mapping.restore_text(run.text)
    for table in document.tables:
        _restore_docx_table(table, mapping)
    for section in document.sections:
        _restore_docx_part(section.header, mapping)
        _restore_docx_part(section.footer, mapping)


def _restore_docx_table(table, mapping: MappingStore) -> None:
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        run.text = mapping.restore_text(run.text)
            for nested_table in cell.tables:
                _restore_docx_table(nested_table, mapping)


def _restore_docx_part(part, mapping: MappingStore) -> None:
    for paragraph in part.paragraphs:
        for run in paragraph.runs:
            if run.text:
                run.text = mapping.restore_text(run.text)
    for table in part.tables:
        _restore_docx_table(table, mapping)


def _iter_docx_header_footer_parts(section):
    for attribute in (
        "header",
        "first_page_header",
        "even_page_header",
        "footer",
        "first_page_footer",
        "even_page_footer",
    ):
        yield getattr(section, attribute)


def _clear_docx_headers_footers(document: Document) -> int:
    cleared_parts = 0
    for section in document.sections:
        for part in _iter_docx_header_footer_parts(section):
            if _header_footer_has_content(part):
                cleared_parts += 1
            try:
                part.is_linked_to_previous = False
            except Exception:
                pass
            _clear_header_footer_part(part)
    return cleared_parts


def _header_footer_has_content(part) -> bool:
    for paragraph in part.paragraphs:
        if paragraph.text.strip():
            return True
        if paragraph.runs:
            return True
    return bool(part.tables)


def _clear_header_footer_part(part) -> None:
    element = part._element
    for child in list(element):
        element.remove(child)
    part.add_paragraph()


def _apply_pdf_header_footer_redactions(document: fitz.Document) -> int:
    candidates_by_page = [_pdf_header_footer_candidates(page) for page in document]
    selected_keys = _pdf_repeated_header_footer_keys(candidates_by_page, len(document))
    removed_count = 0
    for page_index, page in enumerate(document):
        page_removed = False
        for candidate in candidates_by_page[page_index]:
            if candidate["key"] not in selected_keys and not candidate["page_number_like"]:
                continue
            rect = _expand_pdf_rect(page.rect, candidate["rect"])
            page.add_redact_annot(rect, fill=(1, 1, 1), cross_out=False)
            removed_count += 1
            page_removed = True
        if page_removed:
            page.apply_redactions()
    return removed_count


def _pdf_header_footer_candidates(page) -> list[dict[str, object]]:
    page_rect = page.rect
    region_height = page_rect.height * PDF_HEADER_FOOTER_REGION_RATIO
    top_limit = page_rect.y0 + region_height
    bottom_limit = page_rect.y1 - region_height
    candidates: list[dict[str, object]] = []
    text_dict = page.get_text("dict")
    for block in text_dict.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            text = "".join(span.get("text", "") for span in line.get("spans", [])).strip()
            if not text:
                continue
            rect = fitz.Rect(line.get("bbox"))
            if rect.y1 <= top_limit:
                region = "header"
            elif rect.y0 >= bottom_limit:
                region = "footer"
            else:
                continue
            key_text = _normalize_pdf_header_footer_text(text)
            if not key_text:
                continue
            candidates.append(
                {
                    "key": f"{region}:{key_text}",
                    "rect": rect,
                    "page_number_like": _looks_like_pdf_page_number(text),
                }
            )
    return candidates


def _pdf_repeated_header_footer_keys(candidates_by_page: list[list[dict[str, object]]], page_count: int) -> set[str]:
    if page_count < PDF_HEADER_FOOTER_MIN_PAGES:
        return set()
    threshold = min(page_count, max(PDF_HEADER_FOOTER_MIN_PAGES, (page_count + 2) // 3))
    pages_by_key: dict[str, set[int]] = {}
    for page_index, candidates in enumerate(candidates_by_page):
        for candidate in candidates:
            key = str(candidate["key"])
            pages_by_key.setdefault(key, set()).add(page_index)
    return {key for key, pages in pages_by_key.items() if len(pages) >= threshold}


def _normalize_pdf_header_footer_text(text: str) -> str:
    normalized = re.sub(r"\s+", " ", text).strip()
    normalized = re.sub(r"\d+", "#", normalized)
    normalized = normalized.strip("-_/| 第页頁PagepageofOF ")
    return normalized


def _looks_like_pdf_page_number(text: str) -> bool:
    value = re.sub(r"\s+", " ", text).strip()
    return bool(
        re.fullmatch(r"[-–—]?\s*\d+\s*[-–—]?", value)
        or re.fullmatch(r"\d+\s*/\s*\d+", value)
        or re.fullmatch(r"(第\s*)?\d+\s*(页|頁)", value)
        or re.fullmatch(r"Page\s+\d+(\s+of\s+\d+)?", value, flags=re.IGNORECASE)
    )


def _expand_pdf_rect(page_rect: fitz.Rect, rect: fitz.Rect) -> fitz.Rect:
    expanded = fitz.Rect(rect)
    margin = 2.0
    expanded.x0 = max(page_rect.x0, expanded.x0 - margin)
    expanded.y0 = max(page_rect.y0, expanded.y0 - margin)
    expanded.x1 = min(page_rect.x1, expanded.x1 + margin)
    expanded.y1 = min(page_rect.y1, expanded.y1 + margin)
    return expanded


def _replace_pptx_deck(
    presentation,
    mapping: MappingStore,
    custom_terms: Iterable[str],
) -> Counter[str]:
    counts = Counter()
    for slide in presentation.slides:
        counts.update(_replace_pptx_shapes(slide.shapes, mapping, custom_terms))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            counts.update(_replace_pptx_text_frame(slide.notes_slide.notes_text_frame, mapping, custom_terms))
    return counts


def _apply_replacements_pptx_deck(
    presentation,
    replacements: list[ReplacementSpec],
) -> Counter[str]:
    counts = Counter()
    for slide in presentation.slides:
        counts.update(_apply_replacements_pptx_shapes(slide.shapes, replacements))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            counts.update(_apply_replacements_pptx_text_frame(slide.notes_slide.notes_text_frame, replacements))
    return counts


def _restore_pptx_deck(presentation, mapping: MappingStore) -> None:
    for slide in presentation.slides:
        _restore_pptx_shapes(slide.shapes, mapping)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            _restore_pptx_text_frame(slide.notes_slide.notes_text_frame, mapping)


def _replace_pptx_shapes(shapes, mapping: MappingStore, custom_terms: Iterable[str]) -> Counter[str]:
    counts = Counter()
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            counts.update(_replace_pptx_text_frame(shape.text_frame, mapping, custom_terms))
        if getattr(shape, "has_table", False):
            counts.update(_replace_pptx_table(shape.table, mapping, custom_terms))
        if hasattr(shape, "shapes"):
            counts.update(_replace_pptx_shapes(shape.shapes, mapping, custom_terms))
    return counts


def _apply_replacements_pptx_shapes(shapes, replacements: list[ReplacementSpec]) -> Counter[str]:
    counts = Counter()
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            counts.update(_apply_replacements_pptx_text_frame(shape.text_frame, replacements))
        if getattr(shape, "has_table", False):
            counts.update(_apply_replacements_pptx_table(shape.table, replacements))
        if hasattr(shape, "shapes"):
            counts.update(_apply_replacements_pptx_shapes(shape.shapes, replacements))
    return counts


def _restore_pptx_shapes(shapes, mapping: MappingStore) -> None:
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            _restore_pptx_text_frame(shape.text_frame, mapping)
        if getattr(shape, "has_table", False):
            _restore_pptx_table(shape.table, mapping)
        if hasattr(shape, "shapes"):
            _restore_pptx_shapes(shape.shapes, mapping)


def _replace_pptx_table(table, mapping: MappingStore, custom_terms: Iterable[str]) -> Counter[str]:
    counts = Counter()
    for row in table.rows:
        for cell in row.cells:
            counts.update(_replace_pptx_text_frame(cell.text_frame, mapping, custom_terms))
    return counts


def _apply_replacements_pptx_table(table, replacements: list[ReplacementSpec]) -> Counter[str]:
    counts = Counter()
    for row in table.rows:
        for cell in row.cells:
            counts.update(_apply_replacements_pptx_text_frame(cell.text_frame, replacements))
    return counts


def _restore_pptx_table(table, mapping: MappingStore) -> None:
    for row in table.rows:
        for cell in row.cells:
            _restore_pptx_text_frame(cell.text_frame, mapping)


def _replace_pptx_text_frame(text_frame, mapping: MappingStore, custom_terms: Iterable[str]) -> Counter[str]:
    if not text_frame.text:
        return Counter()
    new_text, counts = anonymize_text(text_frame.text, mapping, custom_terms)
    if counts:
        text_frame.text = new_text
    return counts


def _apply_replacements_pptx_text_frame(text_frame, replacements: list[ReplacementSpec]) -> Counter[str]:
    if not text_frame.text:
        return Counter()
    new_text, counts = apply_replacements_text(text_frame.text, replacements)
    if counts:
        text_frame.text = new_text
    return counts


def _restore_pptx_text_frame(text_frame, mapping: MappingStore) -> None:
    if text_frame.text:
        text_frame.text = mapping.restore_text(text_frame.text)


def _restore_pdf_text(input_path: Path, output_dir: Path, mapping: MappingStore) -> tuple[Path, str]:
    output_path = _output_path(input_path, output_dir, "restored", input_path.suffix)
    replacements = [
        ReplacementSpec(
            entity=record.entity,
            prefix=record.prefix,
            value=record.placeholder,
            replacement=record.value,
        )
        for record in mapping.records
        if record.placeholder and record.value
    ]
    if not replacements:
        shutil.copy2(input_path, output_path)
        return output_path, "PDF copied. The mapping file did not contain restore records."
    counts = _restore_pdf_with_replacements(input_path, output_path, replacements)
    if not counts:
        return output_path, "PDF copied. No mapping placeholders were found in the file."
    return output_path, "PDF restored in original layout with mapping placeholders."


def _restore_pdf_with_replacements(
    input_path: Path,
    output_path: Path,
    replacements: Iterable[ReplacementSpec],
) -> Counter[str]:
    _register_pdf_restore_font()
    counts: Counter[str] = Counter()
    overlay_items: dict[int, list[tuple[fitz.Rect, str, float]]] = {}
    replacement_list = sorted(
        _unique_replacements(replacements),
        key=lambda item: len(item.value),
        reverse=True,
    )
    document = fitz.open(str(input_path))
    try:
        for page_index, page in enumerate(document):
            page_items: list[tuple[fitz.Rect, str, float]] = []
            has_redactions = False
            for replacement in replacement_list:
                rectangles = page.search_for(replacement.value)
                if not rectangles:
                    continue
                for rectangle in rectangles:
                    restore_rect, font_size = _pdf_overlay_rect_and_font(page.rect, rectangle, replacement.replacement)
                    page.add_redact_annot(restore_rect, fill=(1, 1, 1), cross_out=False)
                    page_items.append((restore_rect, replacement.replacement, font_size))
                    counts[replacement.entity] += 1
                    has_redactions = True
            if has_redactions:
                page.apply_redactions()
                overlay_items[page_index] = page_items

        if not counts:
            shutil.copy2(input_path, output_path)
            return counts

        _apply_pdf_text_overlays(document, overlay_items)
        document.save(str(output_path), garbage=4, deflate=True, clean=True)
        return counts
    finally:
        document.close()


def _register_pdf_restore_font() -> None:
    try:
        pdfmetrics.getFont(PDF_RESTORE_FONT)
    except KeyError:
        pdfmetrics.registerFont(UnicodeCIDFont(PDF_RESTORE_FONT))


def _pdf_overlay_rect_and_font(page_rect: fitz.Rect, source_rect: fitz.Rect, text: str) -> tuple[fitz.Rect, float]:
    rect = fitz.Rect(source_rect)
    rect.y0 = max(page_rect.y0, rect.y0 - 1)
    rect.y1 = min(page_rect.y1, rect.y1 + 2)
    font_size = min(10.0, max(5.0, rect.height * 0.78))
    available_width = max(1.0, page_rect.x1 - rect.x0 - 2)
    text_width = pdfmetrics.stringWidth(text, PDF_RESTORE_FONT, font_size)
    while text_width > available_width and font_size > 5.0:
        font_size -= 0.5
        text_width = pdfmetrics.stringWidth(text, PDF_RESTORE_FONT, font_size)
    rect.x1 = min(page_rect.x1, max(rect.x1, rect.x0 + text_width + 4))
    return rect, font_size


def _apply_pdf_text_overlays(
    document: fitz.Document,
    overlay_items: dict[int, list[tuple[fitz.Rect, str, float]]],
) -> None:
    for page_index, items in overlay_items.items():
        if not items:
            continue
        page = document[page_index]
        overlay_pdf = _build_pdf_text_overlay(page.rect, items)
        overlay_document = fitz.open("pdf", overlay_pdf)
        try:
            page.show_pdf_page(page.rect, overlay_document, 0, overlay=True)
        finally:
            overlay_document.close()


def _build_pdf_text_overlay(
    page_rect: fitz.Rect,
    items: list[tuple[fitz.Rect, str, float]],
) -> bytes:
    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=(page_rect.width, page_rect.height))
    pdf.setFillColorRGB(0, 0, 0)
    for rect, text, font_size in items:
        baseline = page_rect.height - rect.y1 + max(1.0, (rect.height - font_size) * 0.45)
        pdf.setFont(PDF_RESTORE_FONT, font_size)
        pdf.drawString(rect.x0, baseline, text)
    pdf.save()
    return buffer.getvalue()


def _redact_pdf_with_replacements(
    input_path: Path,
    output_dir: Path,
    replacements: Iterable[ReplacementSpec],
    remove_headers_footers: bool = False,
) -> tuple[Path, Counter[str], str]:
    output_path = _output_path(input_path, output_dir, "desensitized", input_path.suffix)
    _register_pdf_restore_font()
    counts: Counter[str] = Counter()
    overlay_items: dict[int, list[tuple[fitz.Rect, str, float]]] = {}
    replacement_list = sorted(
        _unique_replacements(replacements),
        key=lambda item: len(item.value),
        reverse=True,
    )
    document = fitz.open(str(input_path))
    removed_header_footer_count = 0
    try:
        header_footer_redactions: dict[int, list[fitz.Rect]] = {}
        if remove_headers_footers:
            candidates_by_page = [_pdf_header_footer_candidates(page) for page in document]
            selected_keys = _pdf_repeated_header_footer_keys(candidates_by_page, len(document))
            for page_index, page_candidates in enumerate(candidates_by_page):
                for candidate in page_candidates:
                    if candidate["key"] not in selected_keys and not candidate["page_number_like"]:
                        continue
                    header_footer_redactions.setdefault(page_index, []).append(fitz.Rect(candidate["rect"]))

        for page in document:
            has_redactions = False
            for rectangle in header_footer_redactions.get(page.number, []):
                page.add_redact_annot(_expand_pdf_rect(page.rect, rectangle), fill=(1, 1, 1), cross_out=False)
                removed_header_footer_count += 1
                has_redactions = True
            for replacement in replacement_list:
                if not replacement.value:
                    continue
                rectangles = page.search_for(replacement.value)
                if not rectangles:
                    continue
                for rectangle in rectangles:
                    redact_rect, font_size = _pdf_overlay_rect_and_font(page.rect, rectangle, replacement.replacement)
                    page.add_redact_annot(redact_rect, fill=(1, 1, 1), cross_out=False)
                    overlay_items.setdefault(page.number, []).append((redact_rect, replacement.replacement, font_size))
                    counts[replacement.entity] += 1
                    has_redactions = True
            if has_redactions:
                page.apply_redactions()
        if not counts and not removed_header_footer_count:
            shutil.copy2(input_path, output_path)
        else:
            if overlay_items:
                _apply_pdf_text_overlays(document, overlay_items)
            document.save(str(output_path), garbage=4, deflate=True, clean=True)
    finally:
        document.close()
    message = "PDF redacted in original layout with approved replacements."
    if remove_headers_footers:
        if removed_header_footer_count:
            message += f" Removed {removed_header_footer_count} repeated header/footer item(s)."
        else:
            message += " No repeated PDF header/footer text was detected."
    return output_path, counts, message


def _unique_replacements(replacements: Iterable[ReplacementSpec]) -> list[ReplacementSpec]:
    by_value: dict[str, ReplacementSpec] = {}
    for replacement in replacements:
        if replacement.context_key:
            continue
        if replacement.value and replacement.value not in by_value:
            by_value[replacement.value] = replacement
    return list(by_value.values())


def _detect_table_headers(sheet) -> tuple[int | None, dict[int, str]]:
    best_row_index: int | None = None
    best_headers: dict[int, str] = {}
    best_score = 0
    max_row = min(sheet.max_row or EXCEL_HEADER_SCAN_ROWS, EXCEL_HEADER_SCAN_ROWS)
    for row in sheet.iter_rows(min_row=1, max_row=max_row):
        headers = {cell.column: _cell_text(cell.value) for cell in row if _cell_text(cell.value)}
        if not headers:
            continue
        normalized_headers = {_normalize_header(value) for value in headers.values()}
        score = len(normalized_headers & {_normalize_header(value) for value in PERSON_NAME_HEADERS})
        score += len(normalized_headers & {_normalize_header(value) for value in SUBJECT_ID_HEADERS})
        if score > best_score:
            best_score = score
            best_row_index = row[0].row
            best_headers = headers
    if best_score < 2:
        return None, {}
    return best_row_index, best_headers


def _columns_matching(headers: dict[int, str], names: set[str]) -> list[int]:
    normalized_names = {_normalize_header(value) for value in names}
    return [column for column, header in headers.items() if _normalize_header(header) in normalized_names]


def _first_identifier(
    row_values: dict[int, str],
    headers: dict[int, str],
    id_columns: list[int],
) -> tuple[str, str] | None:
    for column_index in id_columns:
        value = row_values.get(column_index, "").strip()
        if value:
            return headers.get(column_index, f"Column{column_index}"), value
    return None


def _row_context_keys(
    row_values: dict[int, str],
    headers: dict[int, str],
    id_columns: list[int],
) -> list[str]:
    keys: list[str] = []
    for column_index in id_columns:
        value = row_values.get(column_index, "").strip()
        if not value:
            continue
        header = headers.get(column_index, f"Column{column_index}")
        keys.append(_subject_context_key("PERSON", header, value))
    return keys


def _subject_context_key(entity: str, id_header: str, id_value: str) -> str:
    return f"{entity}|{_normalize_header(id_header)}|{id_value.strip()}"


def _normalize_header(value: str) -> str:
    return re.sub(r"[\s:：()（）\[\]【】_-]+", "", str(value)).strip()


def _cell_text(value: object) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _looks_like_person_name(value: str) -> bool:
    value = value.strip()
    return bool(re.fullmatch(r"[\u4e00-\u9fff]{2,4}", value))


def _replacements_by_context(replacements: Iterable[ReplacementSpec]) -> dict[str, list[ReplacementSpec]]:
    scoped: dict[str, list[ReplacementSpec]] = {}
    for replacement in replacements:
        if not replacement.context_key:
            continue
        scoped.setdefault(replacement.context_key, []).append(replacement)
    return scoped


def _scoped_replacements_for_row(
    row_context_keys: list[str],
    scoped_by_context: dict[str, list[ReplacementSpec]],
) -> list[ReplacementSpec]:
    replacements: list[ReplacementSpec] = []
    for context_key in row_context_keys:
        replacements.extend(scoped_by_context.get(context_key, []))
    return replacements


def _apply_scoped_replacements_text(
    text: str,
    replacements: Iterable[ReplacementSpec],
) -> tuple[str, Counter[str]]:
    counts: Counter[str] = Counter()
    new_text = text
    replacement_list = list(replacements)
    codes = {r.replacement for r in replacement_list if r.value and r.value != r.replacement}
    conflicted = _prefix_conflict_codes(codes)
    for replacement in sorted(replacement_list, key=lambda item: len(item.value), reverse=True):
        if not replacement.value or replacement.value == replacement.replacement:
            continue
        found = new_text.count(replacement.value)
        if not found:
            continue
        code = replacement.replacement
        if code in conflicted:
            new_text = new_text.replace(replacement.value, code + ZWSP_MARKER)
        else:
            new_text = new_text.replace(replacement.value, code)
        counts[replacement.entity] += found
    return new_text, counts


def _extract_text_chunks(input_path: Path) -> list[str]:
    suffix = input_path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return [read_text_file(input_path)]
    if suffix in WORD_SUFFIXES:
        document = Document(str(input_path))
        return _iter_docx_text_chunks(document)
    if suffix in EXCEL_SUFFIXES:
        workbook = load_workbook(input_path, data_only=False, read_only=True)
        chunks: list[str] = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str):
                        chunks.append(cell.value)
        workbook.close()
        return chunks
    if suffix in POWERPOINT_SUFFIXES:
        presentation = Presentation(str(input_path))
        return _iter_pptx_text_chunks(presentation)
    if suffix in PDF_SUFFIXES:
        return _extract_pdf_text_chunks(input_path)
    raise DesensitizeError(f"Unsupported file type: {suffix or '(no extension)'}")


def _extract_pdf_text_chunks(input_path: Path) -> list[str]:
    pages: list[str] = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text)
    if not pages:
        raise SkippedFile("该 PDF 没有可提取文字，可能是扫描版。请先转换为文字版 PDF、Word 或文本后再处理。")
    return pages


def _ensure_pdf_has_extractable_text(input_path: Path) -> None:
    _extract_pdf_text_chunks(input_path)


def _iter_docx_text_chunks(document: Document) -> list[str]:
    chunks: list[str] = []
    chunks.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text)
    for table in document.tables:
        chunks.extend(_iter_docx_table_text(table))
    for section in document.sections:
        chunks.extend(paragraph.text for paragraph in section.header.paragraphs if paragraph.text)
        chunks.extend(paragraph.text for paragraph in section.footer.paragraphs if paragraph.text)
        for table in section.header.tables:
            chunks.extend(_iter_docx_table_text(table))
        for table in section.footer.tables:
            chunks.extend(_iter_docx_table_text(table))
    return chunks


def _iter_docx_table_text(table) -> list[str]:
    chunks: list[str] = []
    for row in table.rows:
        for cell in row.cells:
            chunks.extend(paragraph.text for paragraph in cell.paragraphs if paragraph.text)
            for nested_table in cell.tables:
                chunks.extend(_iter_docx_table_text(nested_table))
    return chunks


def _iter_pptx_text_chunks(presentation) -> list[str]:
    chunks: list[str] = []
    for slide in presentation.slides:
        chunks.extend(_iter_pptx_shape_text_chunks(slide.shapes))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            text = slide.notes_slide.notes_text_frame.text
            if text:
                chunks.append(text)
    return chunks


def _iter_pptx_shape_text_chunks(shapes) -> list[str]:
    chunks: list[str] = []
    for shape in shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
            chunks.append(shape.text_frame.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text:
                        chunks.append(cell.text)
        if hasattr(shape, "shapes"):
            chunks.extend(_iter_pptx_shape_text_chunks(shape.shapes))
    return chunks


def _output_path(input_path: Path, output_dir: Path, label: str, suffix: str) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    stem = input_path.stem
    candidate = output_dir / f"{stem}_{label}{suffix}"
    index = 2
    while candidate.exists():
        candidate = output_dir / f"{stem}_{label}_{index}{suffix}"
        index += 1
    return candidate


def copy_unsupported(input_path: Path, output_dir: Path) -> Path:
    output_path = _output_path(input_path, output_dir, "copy", input_path.suffix)
    shutil.copy2(input_path, output_path)
    return output_path
